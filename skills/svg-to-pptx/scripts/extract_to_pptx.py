"""Extract icons from a self-contained SVG and rebuild it as a PowerPoint deck.

Reads an SVG that has base64-inlined `<image>` tags (such as the `final.svg`
produced by the AutoFigure-Edit pipeline) and writes:

- `<output_dir>/assets/icons/<svg_id>.<ext>`  (one PNG/JPG per inlined image)
- `<output_dir>/final_rebuild.pptx`           (single-slide deck rebuilding the figure)

Usage:
    python extract_to_pptx.py [svg_path] [output_dir]

Defaults:
    svg_path   = todo/final.svg
    output_dir = todo

Translated SVG primitives:
- `<image>`    add_picture (icons placed at scaled coords)
- `<rect>`     RECTANGLE / ROUNDED_RECTANGLE
- `<circle>`   OVAL
- `<line>`     CONNECTOR_STRAIGHT
- `<polygon>`  freeform line segments
- `<text>`     textbox (anchor handling, italic via tspan, rotate transform)

Skipped (redraw in PowerPoint):
- `<path>` (most are bezier alignment lines or curved arrows)
- markers / arrowheads
- gradients, drop shadows, dashed strokes
- super/subscript formatting in tspan (kept inline as plain text)
"""
from __future__ import annotations

import argparse
import base64
import re
import sys
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt


SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"

SLIDE_W_INCHES = 13.333  # standard 16:9 widescreen width
EMU_PER_INCH = 914400


def _local(tag) -> str:
    if not isinstance(tag, str):
        return ""
    if tag.startswith("{"):
        return tag.split("}", 1)[1]
    return tag


def _f(elem, attr: str, default: float = 0.0) -> float:
    v = elem.get(attr)
    if v is None or v == "":
        return default
    try:
        return float(v)
    except ValueError:
        return default


def _parse_color(value: Optional[str]) -> Optional[RGBColor]:
    if not value:
        return None
    v = value.strip().lower()
    if v in {"none", "transparent"}:
        return None
    if v.startswith("#"):
        h = v.lstrip("#")
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        if len(h) == 6:
            return RGBColor.from_string(h.upper())
    if v.startswith("rgb("):
        nums = re.findall(r"\d+", v)
        if len(nums) >= 3:
            r, g, b = (int(nums[i]) for i in range(3))
            return RGBColor(r, g, b)
    return None


def _collect_text(elem) -> str:
    """Concatenate text content of <text> with nested <tspan>s."""
    parts: list[str] = []
    if elem.text:
        parts.append(elem.text)
    for child in elem:
        if _local(child.tag) == "tspan":
            if child.text:
                parts.append(child.text)
            for sub in child:
                if sub.text:
                    parts.append(sub.text)
                if sub.tail:
                    parts.append(sub.tail)
        if child.tail:
            parts.append(child.tail)
    return "".join(parts)


def _has_italic(elem) -> bool:
    if elem.get("font-style") == "italic":
        return True
    for child in elem.iter():
        if child.get("font-style") == "italic":
            return True
    return False


def main(svg_path: Path, output_dir: Path) -> None:
    if not svg_path.is_file():
        raise SystemExit(f"missing {svg_path}")
    icons_dir = output_dir / "assets" / "icons"
    icons_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "final_rebuild.pptx"

    tree = etree.parse(str(svg_path))
    root = tree.getroot()

    vb = root.get("viewBox", "")
    parts = vb.split()
    if len(parts) == 4:
        svg_w, svg_h = float(parts[2]), float(parts[3])
    else:
        svg_w = _f(root, "width", 5632.0)
        svg_h = _f(root, "height", 3072.0)

    emu_per_unit = (SLIDE_W_INCHES * EMU_PER_INCH) / svg_w
    pt_per_unit = (SLIDE_W_INCHES * 72.0) / svg_w

    def E(v: float) -> int:
        return int(round(v * emu_per_unit))

    prs = Presentation()
    prs.slide_width = Emu(int(round(svg_w * emu_per_unit)))
    prs.slide_height = Emu(int(round(svg_h * emu_per_unit)))
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    counters = {
        "image": 0, "rect": 0, "text": 0, "polygon": 0,
        "circle": 0, "line": 0, "path_skipped": 0, "other_skipped": 0,
    }

    for elem in root.iter():
        tag = _local(elem.tag)

        if tag == "image":
            href = elem.get(f"{{{XLINK_NS}}}href") or elem.get("href", "")
            if not href.startswith("data:image/"):
                continue
            meta, _, b64 = href.partition(",")
            ext = "png"
            if "jpeg" in meta or "jpg" in meta:
                ext = "jpg"
            elif "webp" in meta:
                ext = "webp"
            data = base64.b64decode(b64)
            iid = elem.get("id") or f"image_{counters['image']:02d}"
            png_path = icons_dir / f"{iid}.{ext}"
            png_path.write_bytes(data)
            counters["image"] += 1
            x = _f(elem, "x"); y = _f(elem, "y")
            w = _f(elem, "width"); h = _f(elem, "height")
            slide.shapes.add_picture(str(png_path), E(x), E(y), E(w), E(h))

        elif tag == "rect":
            x = _f(elem, "x"); y = _f(elem, "y")
            w = _f(elem, "width"); h = _f(elem, "height")
            if w <= 0 or h <= 0:
                continue
            rx = _f(elem, "rx")
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
            shp = slide.shapes.add_shape(shape_type, E(x), E(y), E(w), E(h))
            fill = _parse_color(elem.get("fill"))
            stroke = _parse_color(elem.get("stroke"))
            sw = _f(elem, "stroke-width")
            if fill is not None:
                shp.fill.solid(); shp.fill.fore_color.rgb = fill
            elif elem.get("fill", "").lower() == "none":
                shp.fill.background()
            if stroke is not None:
                shp.line.color.rgb = stroke
                if sw > 0:
                    shp.line.width = Emu(max(E(sw), 6350))
            elif elem.get("stroke", "").lower() == "none":
                shp.line.fill.background()
            counters["rect"] += 1

        elif tag == "circle":
            cx = _f(elem, "cx"); cy = _f(elem, "cy"); r = _f(elem, "r")
            if r <= 0:
                continue
            shp = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, E(cx - r), E(cy - r), E(2 * r), E(2 * r)
            )
            fill = _parse_color(elem.get("fill", "#000000"))
            stroke = _parse_color(elem.get("stroke"))
            sw = _f(elem, "stroke-width")
            if fill is not None:
                shp.fill.solid(); shp.fill.fore_color.rgb = fill
            else:
                shp.fill.background()
            if stroke is not None:
                shp.line.color.rgb = stroke
                if sw > 0:
                    shp.line.width = Emu(max(E(sw), 6350))
            else:
                shp.line.fill.background()
            counters["circle"] += 1

        elif tag == "line":
            x1 = _f(elem, "x1"); y1 = _f(elem, "y1")
            x2 = _f(elem, "x2"); y2 = _f(elem, "y2")
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2),
            )
            stroke = _parse_color(elem.get("stroke", "#000000"))
            sw = _f(elem, "stroke-width", 1.0)
            if stroke is not None:
                conn.line.color.rgb = stroke
            conn.line.width = Emu(max(E(sw), 6350))
            counters["line"] += 1

        elif tag == "polygon":
            pts_str = elem.get("points", "")
            nums = re.findall(r"-?\d+\.?\d*", pts_str)
            pts = [(float(nums[i]), float(nums[i + 1]))
                   for i in range(0, len(nums) - 1, 2)]
            if len(pts) < 3:
                continue
            start_emu_x, start_emu_y = E(pts[0][0]), E(pts[0][1])
            try:
                builder = slide.shapes.build_freeform(start_emu_x, start_emu_y, scale=1.0)
                rest = [(E(p[0]) - start_emu_x, E(p[1]) - start_emu_y) for p in pts[1:]]
                builder.add_line_segments(rest, close=True)
                shp = builder.convert_to_shape()
            except Exception as e:
                print(f"  polygon skipped ({e})")
                continue
            fill = _parse_color(elem.get("fill"))
            stroke = _parse_color(elem.get("stroke"))
            sw = _f(elem, "stroke-width")
            if fill is not None:
                shp.fill.solid(); shp.fill.fore_color.rgb = fill
            else:
                shp.fill.background()
            if stroke is not None:
                shp.line.color.rgb = stroke
                if sw > 0:
                    shp.line.width = Emu(max(E(sw), 6350))
            else:
                shp.line.fill.background()
            counters["polygon"] += 1

        elif tag == "text":
            text = _collect_text(elem).strip()
            if not text:
                continue
            x = _f(elem, "x"); y = _f(elem, "y")
            fs_units = _f(elem, "font-size", 50.0)
            fs_pt = max(6.0, fs_units * pt_per_unit)
            anchor = elem.get("text-anchor", "start")
            fill = _parse_color(elem.get("fill", "#000000")) or RGBColor(0, 0, 0)
            bold = elem.get("font-weight") == "bold"
            italic = _has_italic(elem)

            est_w_units = max(len(text) * fs_units * 0.55, 100.0)
            est_h_units = fs_units * 1.6
            if anchor == "middle":
                tx = x - est_w_units / 2
            elif anchor == "end":
                tx = x - est_w_units
            else:
                tx = x
            ty = y - fs_units

            tb = slide.shapes.add_textbox(E(tx), E(ty), E(est_w_units), E(est_h_units))
            tf = tb.text_frame
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.word_wrap = False
            para = tf.paragraphs[0]
            if anchor == "middle":
                para.alignment = PP_ALIGN.CENTER
            elif anchor == "end":
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = text
            run.font.size = Pt(round(fs_pt, 1))
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = fill

            transform = elem.get("transform", "")
            m = re.search(r"rotate\(\s*(-?\d+\.?\d*)", transform)
            if m:
                tb.rotation = float(m.group(1))
            counters["text"] += 1

        elif tag == "path":
            counters["path_skipped"] += 1

        elif tag in {"svg", "defs", "marker", "tspan", "g"}:
            continue

        else:
            counters["other_skipped"] += 1

    prs.save(str(pptx_path))
    print(f"wrote {pptx_path}")
    print(f"icon files: {icons_dir}")
    print("counters:", counters)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    parser.add_argument("svg_path", nargs="?", default="todo/final.svg")
    parser.add_argument("output_dir", nargs="?", default="todo")
    args = parser.parse_args()
    main(Path(args.svg_path).resolve(), Path(args.output_dir).resolve())
